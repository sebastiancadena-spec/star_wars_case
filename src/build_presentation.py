# ============================ PAQUETES ============================
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================ VARIABLES ============================
OUT = r'c:\Users\sebastian.cadena_kav\Downloads\star_wars_bi_case\output'
PPTX_PATH = f'{OUT}/star_wars_bi_kavak_v2.pptx'

# Colores McKinsey
NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
GOLD  = RGBColor(0xC9, 0xA8, 0x4C)
GRAY  = RGBColor(0x8C, 0x9B, 0xAA)
LGRAY = RGBColor(0xD6, 0xDC, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED   = RGBColor(0xC0, 0x39, 0x2B)
BLACK = RGBColor(0x2C, 0x2C, 0x2C)
BGBOX = RGBColor(0xEF, 0xF2, 0xF6)

W = Inches(13.33)
H = Inches(7.5)

# ============================ FUNCIONES ============================

def nueva_slide(prs, layout_idx=6):
    """
    Agrega una slide en blanco a la presentacion.

    Parametros:
        prs: Objeto Presentation.
        layout_idx (int): Indice del layout (6 = blank).

    Retorna:
        slide: Objeto Slide.
    """
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def fondo(slide, color: RGBColor):
    """
    Establece el color de fondo de una slide.

    Parametros:
        slide: Objeto Slide.
        color (RGBColor): Color de fondo.
    """
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_color: RGBColor = None, line_color: RGBColor = None):
    """
    Agrega un rectangulo a la slide.

    Parametros:
        slide: Objeto Slide.
        left, top, width, height: Posicion y tamano.
        fill_color (RGBColor): Color de relleno opcional.
        line_color (RGBColor): Color de borde opcional.

    Retorna:
        shape: Objeto Shape.
    """
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def texto(slide, text: str, left, top, width, height,
          font_size=12, bold=False, color: RGBColor = None,
          align=PP_ALIGN.LEFT, wrap=True):
    """
    Agrega un cuadro de texto a la slide.

    Parametros:
        slide: Objeto Slide.
        text (str): Texto a mostrar.
        left, top, width, height: Posicion y tamano.
        font_size (int): Tamano de fuente.
        bold (bool): Negrita.
        color (RGBColor): Color de fuente.
        align: Alineacion del parrafo.
        wrap (bool): Word wrap.

    Retorna:
        txBox: Objeto TextBox.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = 'Arial'
    if color:
        run.font.color.rgb = color
    return txBox


def linea_dorada(slide, top):
    """
    Agrega una linea horizontal dorada como separador.

    Parametros:
        slide: Objeto Slide.
        top: Posicion vertical.
    """
    line = slide.shapes.add_connector(1, Inches(0.5), top, Inches(12.83), top)
    line.line.color.rgb = GOLD
    line.line.width = Pt(1.5)


def kpi_box(slide, label: str, value: str, left, top,
            width=Inches(2.8), height=Inches(1.3),
            bg_color: RGBColor = None, val_color: RGBColor = None):
    """
    Agrega un KPI box (valor + etiqueta) a la slide.

    Parametros:
        slide: Objeto Slide.
        label (str): Etiqueta descriptiva.
        value (str): Valor principal.
        left, top: Posicion.
        width, height: Tamano.
        bg_color (RGBColor): Color de fondo.
        val_color (RGBColor): Color del valor.
    """
    bg = bg_color or LGRAY
    vc = val_color or NAVY
    rect(slide, left, top, width, height, fill_color=bg)
    texto(slide, value, left + Inches(0.15), top + Inches(0.12),
          width - Inches(0.3), Inches(0.7),
          font_size=24, bold=True, color=vc, align=PP_ALIGN.CENTER)
    texto(slide, label, left + Inches(0.1), top + Inches(0.8),
          width - Inches(0.2), Inches(0.4),
          font_size=9, bold=False, color=GRAY, align=PP_ALIGN.CENTER)


# ============================ SLIDES ============================

def slide1_titulo(prs):
    """Slide 1: Portada."""
    s = nueva_slide(prs)
    fondo(s, NAVY)

    rect(s, 0, 0, W, Inches(0.08), fill_color=GOLD)
    rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill_color=GOLD)

    texto(s, 'INTELIGENCIA ESTRATEGICA DE\nLA GALAXIA: LOGISTICA Y OPERACIONES',
          Inches(1.0), Inches(1.8), Inches(11), Inches(2.2),
          font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    texto(s, 'Reporte de Inteligencia Estrategica para el Consejo Jedi',
          Inches(1.0), Inches(3.8), Inches(9), Inches(0.6),
          font_size=16, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

    linea_dorada(s, Inches(4.6))

    texto(s, 'Preparado por: Sebastian Cadena  |  Mayo 2026  |  Confidencial',
          Inches(1.0), Inches(4.9), Inches(10), Inches(0.4),
          font_size=10, bold=False, color=GRAY, align=PP_ALIGN.LEFT)

    texto(s, 'MISION: Transformar datos imperiales interceptados en decisiones estrategicas\nsobre eficiencia de flota, talento y optimizacion de inversiones.',
          Inches(1.0), Inches(5.5), Inches(10), Inches(0.9),
          font_size=11, bold=False, color=LGRAY, align=PP_ALIGN.LEFT)


def slide2_flota(prs):
    """Slide 2: Eficiencia de Flota (P1)."""
    s = nueva_slide(prs)
    fondo(s, WHITE)
    rect(s, 0, 0, W, Inches(0.08), fill_color=GOLD)

    texto(s, 'P1 | EFICIENCIA DE FLOTA',
          Inches(0.5), Inches(0.2), Inches(7), Inches(0.35),
          font_size=9, bold=True, color=GOLD)

    texto(s,
          'El Halcon Milenario ofrece el mejor valor: costo por pasajero 6 veces menor\nque cualquier otra nave, pero la evacuacion masiva exige al Executor.',
          Inches(0.5), Inches(0.55), Inches(12.3), Inches(1.0),
          font_size=18, bold=True, color=NAVY)

    linea_dorada(s, Inches(1.6))

    kpi_box(s, 'Naves con datos validos de costo y pasajeros', '4', Inches(0.5), Inches(1.8),
            bg_color=LGRAY, val_color=NAVY)
    kpi_box(s, 'Costo por pasajero - Halcon Milenario', '16,667 creditos', Inches(3.5), Inches(1.8),
            bg_color=GOLD, val_color=WHITE)
    kpi_box(s, 'Capacidad del Executor (pasajeros)', '38,000', Inches(6.5), Inches(1.8),
            bg_color=LGRAY, val_color=NAVY)
    kpi_box(s, 'Costo por pasajero - Executor', '30,088 creditos', Inches(9.5), Inches(1.8),
            bg_color=LGRAY, val_color=NAVY)

    # Tabla ranking
    headers = ['Rank', 'Nave', 'Clase', 'Pasajeros', 'Costo (creditos)', 'Costo por Pasajero', 'Veredicto']
    rows = [
        ['#1', 'Halcon Milenario', 'Carguero Ligero', '6', '100,000', '16,667', 'Mejor Valor'],
        ['#2', 'Executor', 'Destructor Estelar', '38,000', '1,143,350,000', '30,088', 'Max Capacidad'],
        ['#3', 'Transporte Rebelde', 'Transporte Medio', '90', '3,000,000', '33,333', 'Opcion Media'],
    ]

    col_widths = [Inches(0.55), Inches(2.3), Inches(2.0), Inches(1.1), Inches(1.3), Inches(1.5), Inches(1.2)]
    col_lefts = [Inches(0.5)]
    for w in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + w)

    row_top = Inches(3.4)
    for i, (h, w, l) in enumerate(zip(headers, col_widths, col_lefts)):
        rect(s, l, row_top, w - Inches(0.05), Inches(0.38), fill_color=NAVY)
        texto(s, h, l + Inches(0.05), row_top + Inches(0.06), w - Inches(0.1), Inches(0.3),
              font_size=9, bold=True, color=WHITE)

    for r_idx, row in enumerate(rows):
        top = row_top + Inches(0.38) + r_idx * Inches(0.52)
        bg = LGRAY if r_idx % 2 == 0 else WHITE
        highlight = r_idx == 0
        for i, (cell, w, l) in enumerate(zip(row, col_widths, col_lefts)):
            rect(s, l, top, w - Inches(0.05), Inches(0.48),
                 fill_color=GOLD if (highlight and i == 6) else bg)
            texto(s, cell, l + Inches(0.05), top + Inches(0.08), w - Inches(0.1), Inches(0.35),
                  font_size=9, bold=highlight, color=WHITE if (highlight and i == 6) else NAVY)

    rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), fill_color=BGBOX)
    texto(s, 'RECOMENDACION: Para misiones con presupuesto limitado, priorizar el Halcon Milenario '
          '(el costo por pasajero mas bajo de toda la flota). Para evacuacion masiva, el Executor '
          'es la unica nave con capacidad real, pero su precio equivale al 99.97% del presupuesto '
          'total disponible en la flota. No es sostenible depender de una sola nave.',
          Inches(0.65), Inches(6.1), Inches(12.0), Inches(0.7),
          font_size=9, bold=False, color=NAVY)


def slide3_talento(prs):
    """Slide 3: Analisis de Talento (P2)."""
    s = nueva_slide(prs)
    fondo(s, WHITE)
    rect(s, 0, 0, W, Inches(0.08), fill_color=GOLD)

    texto(s, 'P2 | ANALISIS DE TALENTO',
          Inches(0.5), Inches(0.2), Inches(7), Inches(0.35),
          font_size=9, bold=True, color=GOLD)

    texto(s,
          'Tatooine, Ryloth y Naboo producen una tasa de pilotos del 50%, el doble del promedio galactico.\nLos Humanos dominan el cuerpo de pilotos con el 78% de todos los certificados.',
          Inches(0.5), Inches(0.55), Inches(12.3), Inches(1.0),
          font_size=18, bold=True, color=NAVY)

    linea_dorada(s, Inches(1.6))

    kpi_box(s, 'Total personajes analizados', '96', Inches(0.5), Inches(1.8),
            bg_color=LGRAY, val_color=NAVY)
    kpi_box(s, 'Pilotos certificados identificados', '18', Inches(3.5), Inches(1.8),
            bg_color=GOLD, val_color=WHITE)
    kpi_box(s, '1 de cada 5 personajes es piloto', '18.8%', Inches(6.5), Inches(1.8),
            bg_color=LGRAY, val_color=NAVY)
    kpi_box(s, 'Fest y Lothal: sus 2 personajes son pilotos', '100%', Inches(9.5), Inches(1.8),
            bg_color=NAVY, val_color=GOLD)

    # Barras por planeta
    texto(s, 'TASA DE EXITO COMO PILOTO POR PLANETA DE ORIGEN',
          Inches(0.5), Inches(3.25), Inches(6), Inches(0.3),
          font_size=9, bold=True, color=NAVY)

    planet_data = [
        ('Fest', 100), ('Lothal', 100), ('Tatooine', 50),
        ('Ryloth', 50), ('Naboo', 50), ('Alderaan', 50),
        ('Kamino', 50), ('Corellia', 20),
    ]
    bar_max = Inches(4.5)
    for i, (planet, rate) in enumerate(planet_data):
        y = Inches(3.6) + i * Inches(0.38)
        texto(s, planet, Inches(0.5), y + Inches(0.04), Inches(1.2), Inches(0.32),
              font_size=9, color=BLACK)
        bw = bar_max * rate / 100
        bc = GOLD if rate >= 50 else NAVY
        rect(s, Inches(1.75), y + Inches(0.06), bw, Inches(0.24), fill_color=bc)
        texto(s, f'{rate}%', Inches(1.75) + bw + Inches(0.05), y + Inches(0.04),
              Inches(0.5), Inches(0.3), font_size=8, bold=True, color=bc)

    # Tabla por especie
    texto(s, 'PRINCIPALES ESPECIES EN EL CUERPO DE PILOTOS',
          Inches(7.0), Inches(3.25), Inches(5.8), Inches(0.3),
          font_size=9, bold=True, color=NAVY)

    species_data = [
        ('Humano', 14, 60), ("Twi'Lek", 2, 4), ('Wookiee', 1, 1), ('Trandoshan', 1, 1),
    ]
    for i, (sp, pilots, total) in enumerate(species_data):
        y = Inches(3.6) + i * Inches(0.65)
        rate = int(pilots / total * 100)
        rect(s, Inches(7.0), y, Inches(5.7), Inches(0.55),
             fill_color=LGRAY if i % 2 == 0 else WHITE)
        texto(s, sp, Inches(7.1), y + Inches(0.1), Inches(2.0), Inches(0.4),
              font_size=10, color=BLACK)
        texto(s, f'{pilots} pilotos / {total} total = {rate}%',
              Inches(9.2), y + Inches(0.1), Inches(3.3), Inches(0.4),
              font_size=10, bold=True, color=NAVY)

    rect(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5), fill_color=BGBOX)
    texto(s, 'RECOMENDACION: Concentrar el reclutamiento en Tatooine, Ryloth y Naboo. '
          'La brecha de genero es minima (20.6% masculino vs 16.7% femenino) — '
          'el talento esta distribuido equitativamente por genero.',
          Inches(0.65), Inches(6.92), Inches(12.0), Inches(0.4),
          font_size=9, color=NAVY)


def slide4_inversion(prs):
    """Slide 4: Inversion Estrategica (P3)."""
    s = nueva_slide(prs)
    fondo(s, WHITE)
    rect(s, 0, 0, W, Inches(0.08), fill_color=GOLD)

    texto(s, 'P3 | INVERSION ESTRATEGICA',
          Inches(0.5), Inches(0.2), Inches(7), Inches(0.35),
          font_size=9, bold=True, color=GOLD)

    texto(s,
          'Depende del objetivo: el Star Dreadnought gana en escala, el Carguero Ligero gana en presupuesto.\nNunca concentrar la inversion en una sola clase de nave.',
          Inches(0.5), Inches(0.55), Inches(12.3), Inches(1.0),
          font_size=18, bold=True, color=NAVY)

    linea_dorada(s, Inches(1.6))

    # 4 tarjetas por clase
    clases = [
        ('Destructor Estelar', 0.573, '38,000 pasajeros', '1,143 millones creditos',
         'Transporte masivo. Capacidad sin igual\npero costo extremo. Una nave = 99% del presupuesto.', GOLD, WHITE),
        ('Carguero Ligero', 0.500, '6 pasajeros', '100,000 creditos',
         'Mayor agilidad y menor costo por pasajero.\nIdeal para misiones encubiertas y flota distribuida.', NAVY, WHITE),
        ('Corbeta', 0.445, '30 pasajeros', '3.5 millones creditos',
         'Opcion equilibrada. Costo moderado\ny buena velocidad de desplazamiento.', LGRAY, NAVY),
        ('Transporte Medio', 0.300, '90 pasajeros', '3 millones creditos',
         'Buena capacidad a costo moderado,\npero hiperpropulsion lenta limita su alcance.', LGRAY, NAVY),
    ]

    for i, (nombre, score, pax, costo, desc, bg, fg) in enumerate(clases):
        l = Inches(0.5) + i * Inches(3.1)
        rect(s, l, Inches(1.75), Inches(2.95), Inches(4.5), fill_color=bg)

        texto(s, nombre, l + Inches(0.1), Inches(1.85), Inches(2.75), Inches(0.45),
              font_size=11, bold=True, color=fg)
        texto(s, f'Puntaje Estrategico: {score:.2f}', l + Inches(0.1), Inches(2.3),
              Inches(2.75), Inches(0.3), font_size=9, color=fg)

        # Barra de score
        bar_max = Inches(2.5)
        rect(s, l + Inches(0.1), Inches(2.65), bar_max, Inches(0.18),
             fill_color=RGBColor(0xA0, 0xA0, 0xA0))
        rect(s, l + Inches(0.1), Inches(2.65), bar_max * score, Inches(0.18),
             fill_color=WHITE if bg in (NAVY, GOLD) else NAVY)

        texto(s, f'Pasajeros prom.: {pax}', l + Inches(0.1), Inches(2.95), Inches(2.75), Inches(0.28),
              font_size=9, color=fg)
        texto(s, f'Costo prom.: {costo}', l + Inches(0.1), Inches(3.22), Inches(2.75), Inches(0.28),
              font_size=9, color=fg)
        texto(s, desc, l + Inches(0.1), Inches(3.6), Inches(2.75), Inches(0.75),
              font_size=8.5, color=fg)

    rect(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.45), fill_color=BGBOX)
    texto(s, 'Metodologia del puntaje: 50% capacidad de pasajeros + 30% menor costo + 20% mayor velocidad subluminica (cada metrica normalizada de 0 a 1)',
          Inches(0.65), Inches(6.42), Inches(12.0), Inches(0.35),
          font_size=8.5, color=GRAY)

    rect(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5), fill_color=BGBOX)
    texto(s, 'RECOMENDACION: Portafolio balanceado: 1 Dreadnought para operaciones masivas + flota de Cargueros Ligeros '
          'para presencia distribuida. Evitar concentrar el 100% del presupuesto en una sola clase.',
          Inches(0.65), Inches(6.92), Inches(12.0), Inches(0.4),
          font_size=9, color=NAVY)


def slide5_anomalia(prs):
    """Slide 5: Anomalia libre (P4) + Arquitectura."""
    s = nueva_slide(prs)
    fondo(s, WHITE)
    rect(s, 0, 0, W, Inches(0.08), fill_color=RED)

    texto(s, 'P4 | HALLAZGO CRITICO — ANALISIS LIBRE',
          Inches(0.5), Inches(0.2), Inches(9), Inches(0.35),
          font_size=9, bold=True, color=RED)

    texto(s,
          'La Resistencia No Puede Ganar una Guerra de Desgaste:\nLa Flota Puede Transportar Solo el 0.003% de la Poblacion Galactica.',
          Inches(0.5), Inches(0.55), Inches(12.3), Inches(1.0),
          font_size=18, bold=True, color=NAVY)

    linea_dorada(s, Inches(1.6))

    kpi_box(s, 'Habitantes en planetas conocidos de la galaxia', '1.14 Billones', Inches(0.5), Inches(1.8),
            bg_color=RED, val_color=WHITE)
    kpi_box(s, 'Pasajeros que cabe toda la flota (todas las naves juntas)', '38,126', Inches(3.6), Inches(1.8),
            bg_color=LGRAY, val_color=NAVY)
    kpi_box(s, 'Viajes necesarios para evacuar todos los planetas', '29.9 Millones', Inches(6.7), Inches(1.8),
            bg_color=NAVY, val_color=GOLD)
    kpi_box(s, 'Del total de pasajeros de la flota, esta 1 sola nave los tiene', '99.7%', Inches(9.8), Inches(1.8),
            bg_color=RED, val_color=WHITE)

    # Grafico Pareto simplificado
    texto(s, 'CONCENTRACION DE CAPACIDAD (Principio de Pareto)',
          Inches(0.5), Inches(3.25), Inches(6), Inches(0.3),
          font_size=9, bold=True, color=NAVY)

    pareto = [('Executor', 38000), ('Las demas (59 naves)', 126)]
    total = sum(v for _, v in pareto)
    bar_max = Inches(5.5)
    for i, (label, val) in enumerate(pareto):
        y = Inches(3.6) + i * Inches(0.72)
        pct = val / total * 100
        bw = bar_max * pct / 100
        bc = RED if i == 0 else LGRAY
        rect(s, Inches(0.5), y, bw, Inches(0.5), fill_color=bc)
        fc = WHITE if i == 0 else BLACK
        texto(s, f'{label}: {val:,} pax ({pct:.1f}%)',
              Inches(0.6), y + Inches(0.1), Inches(5.5), Inches(0.35),
              font_size=9, bold=(i == 0), color=fc)

    texto(s, '1 sola nave concentra el 99.7% de toda la capacidad de pasajeros de la galaxia.',
          Inches(0.5), Inches(5.1), Inches(6.2), Inches(0.4),
          font_size=10, bold=True, color=RED)

    texto(s, 'POR QUE EL CONSEJO DEBE ACTUAR: Perder el Executor equivale a perder\nla capacidad de evacuar cualquier planeta. Es un punto critico de falla.',
          Inches(0.5), Inches(5.55), Inches(6.2), Inches(0.7),
          font_size=10, color=NAVY)

    # Arquitectura (columna derecha)
    rect(s, Inches(7.0), Inches(3.2), Inches(5.8), Inches(4.15), fill_color=RGBColor(0xF5, 0xF7, 0xFA))
    texto(s, 'ARQUITECTURA UTILIZADA', Inches(7.1), Inches(3.3), Inches(5.6), Inches(0.3),
          font_size=9, bold=True, color=NAVY)

    arch_text = (
        'Los datos crudos en formato parquet fueron procesados con Python en VS Code: '
        'se limpiaron, transformaron y calcularon las metricas de cada pregunta de negocio. '
        'Las tablas resultantes se cargaron en Google Sheets como fuente de datos, '
        'y desde ahi se conectaron a Looker Studio para construir el dashboard publico. '
        'Para produccion a escala, este mismo flujo correria en Google Cloud Platform (GCP) '
        'con Google Drive como almacen de archivos, BigQuery para transformaciones masivas '
        'y Looker Studio como capa de visualizacion final.'
    )
    texto(s, arch_text, Inches(7.1), Inches(3.65), Inches(5.5), Inches(2.1),
          font_size=9, color=BLACK, wrap=True)

    # Diagrama de arquitectura — lo que se hizo
    pasos = ['Archivos\nParquet', 'Python\n(VS Code)', 'Google\nSheets', 'Looker\nStudio']
    for i, paso in enumerate(pasos):
        l = Inches(7.1) + i * Inches(1.4)
        bg = GOLD if i == 2 else (NAVY if i == 3 else LGRAY)
        fc = WHITE if bg in (NAVY, GOLD) else NAVY
        rect(s, l, Inches(5.85), Inches(1.25), Inches(0.6), fill_color=bg)
        texto(s, paso, l + Inches(0.05), Inches(5.9), Inches(1.15), Inches(0.5),
              font_size=7.5, color=fc, align=PP_ALIGN.CENTER)
        if i < 3:
            texto(s, '>', l + Inches(1.3), Inches(5.95), Inches(0.15), Inches(0.35),
                  font_size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    texto(s, 'Lo anterior es lo que se implemento. En produccion se escalaria con GCP y BigQuery.',
          Inches(7.1), Inches(6.55), Inches(5.5), Inches(0.35),
          font_size=7.5, color=GRAY)


# ============================ MAIN ============================

def main():
    """
    Construye la presentacion McKinsey para el Consejo Jedi (5 slides, en espanol).
    """
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide1_titulo(prs)
    slide2_flota(prs)
    slide3_talento(prs)
    slide4_inversion(prs)
    slide5_anomalia(prs)

    os.makedirs(OUT, exist_ok=True)
    prs.save(PPTX_PATH)
    print(f'PPT guardado: {PPTX_PATH}')


if __name__ == '__main__':
    main()
